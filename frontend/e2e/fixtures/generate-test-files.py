#!/usr/bin/env python3
"""
Script pour générer les fichiers de test (Excel, PDF, Word, PowerPoint)
Usage: python generate-test-files.py
"""

import os
from datetime import datetime, timedelta
from pathlib import Path

# Determine fixture directory
FIXTURE_DIR = Path(__file__).parent


def generate_excel_file():
    """Génère un fichier Excel de test avec données Supply Chain."""
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill
    except ImportError:
        print("⚠️  openpyxl non installé. Installer avec: pip install openpyxl")
        return

    wb = openpyxl.Workbook()

    # Feuille 1: Stocks
    ws_stocks = wb.active
    ws_stocks.title = "Stocks"

    # Headers avec style
    headers_stocks = ["SKU", "Produit", "Stock Actuel", "Stock Sécurité", "Statut", "Emplacement"]
    ws_stocks.append(headers_stocks)
    for cell in ws_stocks[1]:
        cell.font = Font(bold=True)
        cell.fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")

    # Données avec cas de test spécifiques
    stocks_data = [
        ["SKU001", "Widget A", 150, 100, "OK", "Entrepôt A"],
        ["SKU002", "Widget B", -50, 75, "RUPTURE", "Entrepôt A"],  # Stock négatif (alerte)
        ["SKU003", "Gadget C", 300, 200, "OK", "Entrepôt B"],
        ["SKU004", "Gadget D", 80, 100, "ALERTE", "Entrepôt B"],  # Stock < sécurité
        ["SKU005", "Tool E", 200, 150, "OK", "Entrepôt C"],
    ]

    for row in stocks_data:
        ws_stocks.append(row)

    # Feuille 2: Commandes
    ws_orders = wb.create_sheet("Commandes")
    headers_orders = ["N° Commande", "SKU", "Quantité", "Date Commande", "Date Livraison Prévue", "Fournisseur", "Statut"]
    ws_orders.append(headers_orders)
    for cell in ws_orders[1]:
        cell.font = Font(bold=True)
        cell.fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")

    # Données avec dates calculables pour lead time
    base_date = datetime(2024, 1, 15)
    orders_data = [
        ["CMD001", "SKU001", 200, base_date, base_date + timedelta(days=7), "Fournisseur Alpha", "En cours"],
        ["CMD002", "SKU002", 150, base_date + timedelta(days=3), base_date + timedelta(days=10), "Fournisseur Beta", "En retard"],
        ["CMD003", "SKU003", 250, base_date - timedelta(days=5), base_date + timedelta(days=2), "Fournisseur Gamma", "Livrée"],
        ["CMD004", "SKU004", 120, base_date + timedelta(days=5), base_date + timedelta(days=15), "Fournisseur Alpha", "Planifiée"],
        ["CMD005", "SKU005", 300, base_date - timedelta(days=2), base_date + timedelta(days=28), "Fournisseur Delta", "En cours"],  # Lead time long
    ]

    for row in orders_data:
        ws_orders.append(row)

    # Feuille 3: Fournisseurs
    ws_suppliers = wb.create_sheet("Fournisseurs")
    headers_suppliers = ["Code", "Nom", "Lead Time Moyen (jours)", "Fiabilité (%)", "Pays"]
    ws_suppliers.append(headers_suppliers)
    for cell in ws_suppliers[1]:
        cell.font = Font(bold=True)
        cell.fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")

    suppliers_data = [
        ["SUPP001", "Fournisseur Alpha", 7, 95, "France"],
        ["SUPP002", "Fournisseur Beta", 10, 88, "Allemagne"],
        ["SUPP003", "Fournisseur Gamma", 5, 98, "France"],
        ["SUPP004", "Fournisseur Delta", 30, 75, "Chine"],  # Lead time long, fiabilité basse
    ]

    for row in suppliers_data:
        ws_suppliers.append(row)

    # Sauvegarder
    output_path = FIXTURE_DIR / "test-production.xlsx"
    wb.save(output_path)
    print(f"✅ Fichier Excel créé: {output_path}")


def generate_pdf_file():
    """Génère un fichier PDF de test."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib import colors
    except ImportError:
        print("⚠️  reportlab non installé. Installer avec: pip install reportlab")
        return

    output_path = FIXTURE_DIR / "test-report.pdf"
    doc = SimpleDocTemplate(str(output_path), pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    # Titre
    title = Paragraph("Rapport d'Analyse Supply Chain - Janvier 2024", styles['Title'])
    story.append(title)
    story.append(Spacer(1, 12))

    # Contenu
    content = """
    <b>Résumé Exécutif</b><br/>
    Ce rapport présente une analyse des stocks et des commandes pour le mois de janvier 2024.
    Plusieurs alertes ont été identifiées concernant les ruptures de stock et les retards de livraison.
    <br/><br/>
    <b>Points Clés:</b><br/>
    • Stock négatif détecté pour SKU002 (Widget B): -50 unités<br/>
    • 3 commandes en retard identifiées<br/>
    • Lead time moyen: 12 jours<br/>
    • Taux de service: 87%<br/>
    <br/>
    <b>Recommandations:</b><br/>
    1. Augmenter le stock de sécurité pour les produits critiques<br/>
    2. Renégocier les délais avec le Fournisseur Delta<br/>
    3. Mettre en place un système d'alerte automatique<br/>
    """

    story.append(Paragraph(content, styles['Normal']))
    story.append(Spacer(1, 12))

    # Tableau
    data = [
        ['SKU', 'Produit', 'Stock', 'Statut'],
        ['SKU001', 'Widget A', '150', 'OK'],
        ['SKU002', 'Widget B', '-50', 'RUPTURE'],
        ['SKU003', 'Gadget C', '300', 'OK'],
    ]

    table = Table(data)
    table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
    ]))

    story.append(table)

    doc.build(story)
    print(f"✅ Fichier PDF créé: {output_path}")


def generate_word_file():
    """Génère un fichier Word de test."""
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
    except ImportError:
        print("⚠️  python-docx non installé. Installer avec: pip install python-docx")
        return

    doc = Document()

    # Titre
    title = doc.add_heading('Procédure de Gestion des Stocks', level=1)

    # Introduction
    doc.add_paragraph(
        "Ce document décrit la procédure standard pour la gestion des stocks dans notre entrepôt. "
        "Il est essentiel de suivre ces étapes pour maintenir un niveau de service optimal."
    )

    # Section 1
    doc.add_heading('1. Contrôle des Stocks', level=2)
    doc.add_paragraph(
        "Vérifier quotidiennement les niveaux de stock. "
        "Toute valeur négative doit être signalée immédiatement au responsable. "
        "Le stock de sécurité doit être maintenu à 100 unités minimum pour les produits critiques."
    )

    # Section 2
    doc.add_heading('2. Gestion des Commandes', level=2)
    doc.add_paragraph(
        "Les commandes doivent être passées dès que le stock atteint le point de commande. "
        "Le lead time moyen avec nos fournisseurs est de 7-10 jours. "
        "Prévoir un délai supplémentaire pour les fournisseurs internationaux (30 jours)."
    )

    # Liste à puces
    doc.add_paragraph('Points de vigilance:', style='List Bullet')
    doc.add_paragraph('Vérifier la disponibilité des fournisseurs', style='List Bullet')
    doc.add_paragraph('Confirmer les délais de livraison', style='List Bullet')
    doc.add_paragraph('Suivre les expéditions en temps réel', style='List Bullet')

    output_path = FIXTURE_DIR / "test-document.docx"
    doc.save(output_path)
    print(f"✅ Fichier Word créé: {output_path}")


def generate_powerpoint_file():
    """Génère un fichier PowerPoint de test."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("⚠️  python-pptx non installé. Installer avec: pip install python-pptx")
        return

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Titre
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Revue Supply Chain Q1 2024"
    subtitle.text = "Analyse des Stocks et Performance Fournisseurs"

    # Slide 2: Analyse des stocks
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Analyse des Stocks"
    tf = body_shape.text_frame
    tf.text = "Points Clés:"

    p = tf.add_paragraph()
    p.text = "Stock négatif détecté: SKU002 (-50 unités)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "4 produits sous le seuil de sécurité"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Taux de rotation: 6.2x par an"
    p.level = 1

    # Slide 3: Performance fournisseurs
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Performance Fournisseurs"
    tf = body_shape.text_frame
    tf.text = "Lead Times:"

    p = tf.add_paragraph()
    p.text = "Fournisseur Alpha: 7 jours (95% fiabilité)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Fournisseur Beta: 10 jours (88% fiabilité)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Fournisseur Delta: 30 jours (75% fiabilité) ⚠️"
    p.level = 1

    output_path = FIXTURE_DIR / "test-presentation.pptx"
    prs.save(output_path)
    print(f"✅ Fichier PowerPoint créé: {output_path}")


if __name__ == "__main__":
    print("🔨 Génération des fichiers de test...")
    print()

    generate_excel_file()
    generate_pdf_file()
    generate_word_file()
    generate_powerpoint_file()

    print()
    print("✅ Tous les fichiers de test ont été générés!")
    print(f"📁 Emplacement: {FIXTURE_DIR}")
