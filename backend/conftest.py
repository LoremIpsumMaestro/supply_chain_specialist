"""Pytest configuration and global fixtures."""

import io
import json
import os
from pathlib import Path
from typing import Generator
from uuid import uuid4

import pytest
from openpyxl import Workbook
from PyPDF2 import PdfWriter
from reportlab.pdfgen import canvas
from docx import Document as DocxDocument
from pptx import Presentation

from backend.models.file import FileType
from backend.services.document_parser import DocumentChunk


# ============================================================================
# Test Data Fixtures
# ============================================================================

@pytest.fixture
def test_user_id() -> str:
    """Generate a test user ID."""
    return str(uuid4())


@pytest.fixture
def test_file_id() -> str:
    """Generate a test file ID."""
    return str(uuid4())


# ============================================================================
# Document Fixtures (Binary)
# ============================================================================

@pytest.fixture
def sample_excel_bytes() -> bytes:
    """
    Create a sample Excel file with supply chain data.

    Contains:
    - Sheet 'Stocks' with product inventory data
    - Sheet 'Orders' with order information
    """
    wb = Workbook()

    # Create Stocks sheet
    ws_stocks = wb.active
    ws_stocks.title = "Stocks"
    ws_stocks.append(["Product", "Stock", "Minimum", "Status"])
    ws_stocks.append(["Product A", 150, 100, "OK"])
    ws_stocks.append(["Product B", -50, 100, "Critical"])  # Negative stock alert!
    ws_stocks.append(["Product C", 200, 100, "OK"])
    ws_stocks.append(["Product D", 0, 50, "Empty"])

    # Create Orders sheet
    ws_orders = wb.create_sheet("Orders")
    ws_orders.append(["Order ID", "Product", "Quantity", "Status"])
    ws_orders.append(["ORD001", "Product A", 50, "Shipped"])
    ws_orders.append(["ORD002", "Product B", -10, "Error"])  # Negative qty alert!
    ws_orders.append(["ORD003", "Product C", 75, "Processing"])

    # Save to bytes
    buffer = io.BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    return buffer.read()


@pytest.fixture
def sample_pdf_bytes() -> bytes:
    """Create a sample PDF file with supply chain content."""
    buffer = io.BytesIO()
    pdf_canvas = canvas.Canvas(buffer)

    # Page 1
    pdf_canvas.drawString(100, 800, "Supply Chain Report")
    pdf_canvas.drawString(100, 780, "Quarter 1 - 2024")
    pdf_canvas.drawString(100, 750, "Total inventory value: $1,250,000")
    pdf_canvas.drawString(100, 730, "Lead time average: 15 days")
    pdf_canvas.showPage()

    # Page 2
    pdf_canvas.drawString(100, 800, "Critical Alerts")
    pdf_canvas.drawString(100, 780, "- Product B: Negative stock detected (-50 units)")
    pdf_canvas.drawString(100, 760, "- Supplier X: Lead time 95 days (exceeds threshold)")
    pdf_canvas.showPage()

    pdf_canvas.save()
    buffer.seek(0)
    return buffer.read()


@pytest.fixture
def sample_word_bytes() -> bytes:
    """Create a sample Word document."""
    doc = DocxDocument()
    doc.add_heading('Supply Chain Analysis', 0)
    doc.add_paragraph('This document contains supply chain analysis.')
    doc.add_paragraph('Key findings:')
    doc.add_paragraph('- Inventory levels are stable', style='List Bullet')
    doc.add_paragraph('- Lead times are within acceptable range', style='List Bullet')

    # Add table
    table = doc.add_table(rows=3, cols=3)
    table.cell(0, 0).text = 'Product'
    table.cell(0, 1).text = 'Stock'
    table.cell(0, 2).text = 'Status'
    table.cell(1, 0).text = 'Product A'
    table.cell(1, 1).text = '150'
    table.cell(1, 2).text = 'OK'
    table.cell(2, 0).text = 'Product B'
    table.cell(2, 1).text = '-50'
    table.cell(2, 2).text = 'Critical'

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.read()


@pytest.fixture
def sample_powerpoint_bytes() -> bytes:
    """Create a sample PowerPoint presentation."""
    prs = Presentation()

    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Supply Chain Dashboard"
    slide1.placeholders[1].text = "Q1 2024 Performance"

    # Slide 2: Content
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Inventory Status"
    body = slide2.shapes.placeholders[1].text_frame
    body.text = "Current inventory levels are stable"
    p = body.add_paragraph()
    p.text = "Critical: Product B has negative stock"

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()


@pytest.fixture
def sample_csv_bytes() -> bytes:
    """Create a sample CSV file."""
    csv_content = """Product,Stock,Minimum,Lead Time
Product A,150,100,15
Product B,-50,100,20
Product C,200,100,12
Product D,0,50,18
"""
    return csv_content.encode('utf-8')


@pytest.fixture
def sample_text_bytes() -> bytes:
    """Create a sample text file."""
    text_content = """Supply Chain Report - Q1 2024

Inventory Summary:
- Total products: 50
- Total stock value: $1.2M
- Average lead time: 16 days

Critical Issues:
- Product B: Negative stock (-50 units)
- Supplier delays: 3 orders delayed over 90 days

Recommendations:
1. Investigate Product B inventory discrepancy
2. Review supplier contracts for lead time guarantees
3. Increase safety stock for high-demand items
"""
    return text_content.encode('utf-8')


# ============================================================================
# Mock Service Fixtures
# ============================================================================

@pytest.fixture
def mock_redis(mocker):
    """Mock Redis client."""
    mock = mocker.Mock()
    mock.get.return_value = None
    mock.setex.return_value = True
    return mock


@pytest.fixture
def mock_typesense(mocker):
    """Mock TypeSense client."""
    mock = mocker.Mock()

    # Mock collection operations
    mock.collections.__getitem__.return_value.retrieve.return_value = {
        'name': 'document_chunks',
        'fields': []
    }

    mock.collections.__getitem__.return_value.documents.import_.return_value = []
    mock.collections.__getitem__.return_value.documents.search.return_value = {
        'hits': []
    }

    return mock


@pytest.fixture
def mock_ollama_embeddings(mocker):
    """Mock Ollama API for embeddings."""
    import requests

    mock_response = mocker.Mock()
    mock_response.json.return_value = {
        'embedding': [0.1] * 768  # Mock 768-dimensional embedding
    }
    mock_response.raise_for_status = mocker.Mock()

    mocker.patch.object(requests, 'post', return_value=mock_response)
    return mock_response


# ============================================================================
# Document Chunk Fixtures
# ============================================================================

@pytest.fixture
def sample_excel_chunks() -> list:
    """Sample chunks extracted from Excel file."""
    return [
        DocumentChunk(
            content="Stock: 150",
            metadata={
                "filename": "test.xlsx",
                "file_type": "excel",
                "sheet_name": "Stocks",
                "cell_ref": "B2",
                "row": 2,
                "column": 2,
                "value": "150",
                "column_header": "Stock",
            },
            file_type=FileType.EXCEL,
        ),
        DocumentChunk(
            content="Stock: -50",
            metadata={
                "filename": "test.xlsx",
                "file_type": "excel",
                "sheet_name": "Stocks",
                "cell_ref": "B3",
                "row": 3,
                "column": 2,
                "value": "-50",
                "column_header": "Stock",
            },
            file_type=FileType.EXCEL,
        ),
    ]


@pytest.fixture
def sample_chunks_with_alerts() -> list:
    """Sample chunks that should trigger alerts."""
    return [
        # Negative stock
        DocumentChunk(
            content="Stock: -50",
            metadata={
                "filename": "inventory.xlsx",
                "file_type": "excel",
                "sheet_name": "Products",
                "cell_ref": "C12",
                "row": 12,
                "column": 3,
                "value": "-50",
                "column_header": "Stock",
            },
            file_type=FileType.EXCEL,
        ),
        # Negative quantity
        DocumentChunk(
            content="Quantity: -10",
            metadata={
                "filename": "orders.xlsx",
                "file_type": "excel",
                "sheet_name": "Orders",
                "cell_ref": "D5",
                "row": 5,
                "column": 4,
                "value": "-10",
                "column_header": "Quantity",
            },
            file_type=FileType.EXCEL,
        ),
        # Long lead time
        DocumentChunk(
            content="Lead time for Product X is 120 days",
            metadata={
                "filename": "suppliers.csv",
                "file_type": "csv",
                "row_number": 10,
            },
            file_type=FileType.CSV,
        ),
    ]


# ============================================================================
# File Path Fixtures
# ============================================================================

@pytest.fixture
def fixtures_dir() -> Path:
    """Get the fixtures directory path."""
    return Path(__file__).parent / "tests" / "fixtures"


@pytest.fixture
def sample_excel_file(fixtures_dir, sample_excel_bytes) -> Path:
    """Create a sample Excel file in the fixtures directory."""
    fixtures_dir.mkdir(parents=True, exist_ok=True)
    file_path = fixtures_dir / "test_inventory.xlsx"
    file_path.write_bytes(sample_excel_bytes)
    yield file_path
    # Cleanup
    if file_path.exists():
        file_path.unlink()


@pytest.fixture
def sample_pdf_file(fixtures_dir, sample_pdf_bytes) -> Path:
    """Create a sample PDF file in the fixtures directory."""
    fixtures_dir.mkdir(parents=True, exist_ok=True)
    file_path = fixtures_dir / "test_report.pdf"
    file_path.write_bytes(sample_pdf_bytes)
    yield file_path
    # Cleanup
    if file_path.exists():
        file_path.unlink()
