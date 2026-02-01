"""Document parsers for extracting structured text with metadata."""

import io
import logging
from abc import ABC, abstractmethod
from dataclasses import dataclass
from datetime import datetime
from typing import List, Optional, Dict, Any

import pandas as pd
from openpyxl import load_workbook
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

from backend.models.file import FileType


logger = logging.getLogger(__name__)


@dataclass
class DocumentChunk:
    """Structured chunk of document content with metadata for RAG."""
    content: str
    metadata: Dict[str, Any]
    file_type: FileType

    def __post_init__(self):
        """Validate chunk after initialization."""
        if not self.content or not self.content.strip():
            raise ValueError("Chunk content cannot be empty")
        if len(self.content) > 10000:  # Max ~2500 tokens
            logger.warning(f"Chunk content very long: {len(self.content)} characters")


class DocumentParser(ABC):
    """Abstract base class for document parsers."""

    @abstractmethod
    def parse(self, file_bytes: bytes, filename: str) -> List[DocumentChunk]:
        """Parse document and return chunks with metadata."""
        pass


class ExcelParser(DocumentParser):
    """Parser for Excel files (.xlsx, .xls) with cell-level granularity."""

    def parse(self, file_bytes: bytes, filename: str) -> List[DocumentChunk]:
        """
        Parse Excel file cell-by-cell.

        Returns one chunk per non-empty cell with metadata:
        - filename, file_type, sheet_name, cell_ref, row, column, value
        """
        chunks = []

        try:
            # Load workbook from bytes
            workbook = load_workbook(io.BytesIO(file_bytes), data_only=True)
            upload_date = datetime.utcnow().isoformat()

            # Import temporal service for temporal analysis
            from backend.services.temporal_service import temporal_service

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]

                # Get column headers from first row (if they exist)
                headers = {}
                for col_idx, cell in enumerate(sheet[1], start=1):
                    if cell.value:
                        headers[col_idx] = str(cell.value)

                # Convert sheet to DataFrame for temporal analysis
                data = []
                for row in sheet.iter_rows(min_row=2, values_only=True):  # Skip header
                    data.append(row)

                if data and headers:
                    df = pd.DataFrame(data, columns=list(headers.values()))

                    # Detect temporal columns
                    temporal_cols = temporal_service.detect_temporal_columns(df, df.columns.tolist())
                    logger.info(f"Detected temporal columns in {sheet_name}: {temporal_cols}")

                    # Identify lead time pairs
                    lead_time_pairs = temporal_service.identify_lead_time_pairs(temporal_cols)
                    lead_time_stats_map = {}

                    # Calculate lead times for each pair
                    for start_col, end_col in lead_time_pairs:
                        stats = temporal_service.calculate_lead_times(df, start_col, end_col)
                        if stats:
                            lead_time_stats_map[f"{start_col}_to_{end_col}"] = stats
                            logger.info(f"Calculated lead times: {start_col} → {end_col}")

                    # Convert DataFrame dates to dict for row lookup
                    temporal_data_by_row = {}
                    if temporal_cols:
                        for row_idx in range(len(df)):
                            temporal_data_by_row[row_idx] = {}
                            for col in temporal_cols:
                                val = df.iloc[row_idx][col]
                                if pd.notna(val):
                                    try:
                                        # Try to convert to datetime and format
                                        dt = pd.to_datetime(val)
                                        temporal_data_by_row[row_idx][col] = dt.strftime('%Y-%m-%d')
                                    except Exception:
                                        temporal_data_by_row[row_idx][col] = str(val)

                # Iterate through all cells
                for row in sheet.iter_rows(min_row=1):
                    for cell in row:
                        if cell.value is not None and str(cell.value).strip():
                            # Build contextual content
                            value = str(cell.value)
                            column_header = headers.get(cell.column, f"Column {cell.column}")

                            # Content includes context for better RAG
                            content = f"{column_header}: {value}"

                            # Base metadata
                            metadata = {
                                "filename": filename,
                                "file_type": "excel",
                                "sheet_name": sheet_name,
                                "cell_ref": cell.coordinate,  # e.g., "C12"
                                "row": cell.row,
                                "column": cell.column,
                                "value": value,
                                "column_header": column_header,
                                "upload_date": upload_date,
                            }

                            # Add temporal context if this row has temporal data
                            row_idx = cell.row - 2  # Adjust for 0-indexed df (and skip header)
                            if row_idx >= 0 and row_idx in temporal_data_by_row:
                                temporal_context = temporal_data_by_row[row_idx]
                                if temporal_context:
                                    metadata["temporal_context"] = temporal_context

                            chunk = DocumentChunk(
                                content=content,
                                metadata=metadata,
                                file_type=FileType.EXCEL,
                            )
                            chunks.append(chunk)

            logger.info(f"Parsed Excel file: {filename} - {len(chunks)} chunks")
            return chunks

        except Exception as e:
            logger.error(f"Error parsing Excel file {filename}: {e}")
            raise


class PDFParser(DocumentParser):
    """Parser for PDF files with page-level chunking."""

    def parse(self, file_bytes: bytes, filename: str) -> List[DocumentChunk]:
        """
        Parse PDF file page-by-page.

        Returns chunks with metadata:
        - filename, file_type, page, chunk_index (if split)
        """
        chunks = []
        upload_date = datetime.utcnow().isoformat()

        try:
            pdf_reader = PdfReader(io.BytesIO(file_bytes))

            for page_num, page in enumerate(pdf_reader.pages, start=1):
                text = page.extract_text()

                if text and text.strip():
                    # Split long pages into multiple chunks (< 1000 tokens ≈ 4000 chars)
                    if len(text) > 4000:
                        # Split by paragraphs
                        paragraphs = text.split('\n\n')
                        current_chunk = ""
                        chunk_index = 0

                        for para in paragraphs:
                            if len(current_chunk) + len(para) < 4000:
                                current_chunk += para + "\n\n"
                            else:
                                if current_chunk:
                                    chunk = DocumentChunk(
                                        content=current_chunk.strip(),
                                        metadata={
                                            "filename": filename,
                                            "file_type": "pdf",
                                            "page": page_num,
                                            "chunk_index": chunk_index,
                                            "upload_date": upload_date,
                                        },
                                        file_type=FileType.PDF,
                                    )
                                    chunks.append(chunk)
                                    chunk_index += 1
                                current_chunk = para + "\n\n"

                        # Add remaining chunk
                        if current_chunk.strip():
                            chunk = DocumentChunk(
                                content=current_chunk.strip(),
                                metadata={
                                    "filename": filename,
                                    "file_type": "pdf",
                                    "page": page_num,
                                    "chunk_index": chunk_index,
                                },
                                file_type=FileType.PDF,
                            )
                            chunks.append(chunk)
                    else:
                        # Page fits in one chunk
                        chunk = DocumentChunk(
                            content=text.strip(),
                            metadata={
                                "filename": filename,
                                "file_type": "pdf",
                                "page": page_num,
                                "chunk_index": 0,
                                "upload_date": upload_date,
                            },
                            file_type=FileType.PDF,
                        )
                        chunks.append(chunk)

            logger.info(f"Parsed PDF file: {filename} - {len(chunks)} chunks")
            return chunks

        except Exception as e:
            logger.error(f"Error parsing PDF file {filename}: {e}")
            raise


class WordParser(DocumentParser):
    """Parser for Word documents (.docx) with paragraph-level chunking."""

    def parse(self, file_bytes: bytes, filename: str) -> List[DocumentChunk]:
        """
        Parse Word file paragraph-by-paragraph.

        Returns chunks with metadata:
        - filename, file_type, paragraph_index
        """
        chunks = []
        upload_date = datetime.utcnow().isoformat()

        try:
            doc = DocxDocument(io.BytesIO(file_bytes))

            for para_idx, paragraph in enumerate(doc.paragraphs):
                text = paragraph.text.strip()

                if text:
                    chunk = DocumentChunk(
                        content=text,
                        metadata={
                            "filename": filename,
                            "file_type": "word",
                            "paragraph_index": para_idx,
                            "upload_date": upload_date,
                        },
                        file_type=FileType.WORD,
                    )
                    chunks.append(chunk)

            # Also parse tables if any
            for table_idx, table in enumerate(doc.tables):
                for row_idx, row in enumerate(table.rows):
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        chunk = DocumentChunk(
                            content=row_text,
                            metadata={
                                "filename": filename,
                                "file_type": "word",
                                "table_index": table_idx,
                                "row_index": row_idx,
                                "upload_date": upload_date,
                            },
                            file_type=FileType.WORD,
                        )
                        chunks.append(chunk)

            logger.info(f"Parsed Word file: {filename} - {len(chunks)} chunks")
            return chunks

        except Exception as e:
            logger.error(f"Error parsing Word file {filename}: {e}")
            raise


class PowerPointParser(DocumentParser):
    """Parser for PowerPoint presentations (.pptx) with slide-level chunking."""

    def parse(self, file_bytes: bytes, filename: str) -> List[DocumentChunk]:
        """
        Parse PowerPoint file slide-by-slide.

        Returns chunks with metadata:
        - filename, file_type, slide_number, content_type (title/body/notes)
        """
        chunks = []
        upload_date = datetime.utcnow().isoformat()

        try:
            prs = Presentation(io.BytesIO(file_bytes))

            for slide_num, slide in enumerate(prs.slides, start=1):
                # Extract title
                if slide.shapes.title:
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        chunk = DocumentChunk(
                            content=title_text,
                            metadata={
                                "filename": filename,
                                "file_type": "powerpoint",
                                "slide_number": slide_num,
                                "content_type": "title",
                                "upload_date": upload_date,
                            },
                            file_type=FileType.POWERPOINT,
                        )
                        chunks.append(chunk)

                # Extract body text
                body_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if shape != slide.shapes.title:  # Skip title
                            body_texts.append(shape.text.strip())

                if body_texts:
                    body_content = "\n".join(body_texts)
                    chunk = DocumentChunk(
                        content=body_content,
                        metadata={
                            "filename": filename,
                            "file_type": "powerpoint",
                            "slide_number": slide_num,
                            "content_type": "body",
                            "upload_date": upload_date,
                        },
                        file_type=FileType.POWERPOINT,
                    )
                    chunks.append(chunk)

                # Extract notes if any
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        chunk = DocumentChunk(
                            content=notes_text,
                            metadata={
                                "filename": filename,
                                "file_type": "powerpoint",
                                "slide_number": slide_num,
                                "content_type": "notes",
                                "upload_date": upload_date,
                            },
                            file_type=FileType.POWERPOINT,
                        )
                        chunks.append(chunk)

            logger.info(f"Parsed PowerPoint file: {filename} - {len(chunks)} chunks")
            return chunks

        except Exception as e:
            logger.error(f"Error parsing PowerPoint file {filename}: {e}")
            raise


class CSVParser(DocumentParser):
    """Parser for CSV files with row-level chunking."""

    def parse(self, file_bytes: bytes, filename: str) -> List[DocumentChunk]:
        """
        Parse CSV file row-by-row.

        Returns chunks with metadata:
        - filename, file_type, row_number, column_headers
        """
        chunks = []

        try:
            # Try to read CSV with pandas
            df = pd.read_csv(io.BytesIO(file_bytes))
            upload_date = datetime.utcnow().isoformat()

            column_headers = df.columns.tolist()

            # Import temporal service for temporal analysis
            from backend.services.temporal_service import temporal_service

            # Detect temporal columns
            temporal_cols = temporal_service.detect_temporal_columns(df, df.columns.tolist())
            logger.info(f"Detected temporal columns in CSV: {temporal_cols}")

            # Calculate temporal data for each row
            temporal_data_by_row = {}
            if temporal_cols:
                for row_idx in range(len(df)):
                    temporal_data_by_row[row_idx] = {}
                    for col in temporal_cols:
                        val = df.iloc[row_idx][col]
                        if pd.notna(val):
                            try:
                                dt = pd.to_datetime(val)
                                temporal_data_by_row[row_idx][col] = dt.strftime('%Y-%m-%d')
                            except Exception:
                                temporal_data_by_row[row_idx][col] = str(val)

            for row_idx, row in df.iterrows():
                # Build content with column headers for context
                row_content = []
                for col in df.columns:
                    value = row[col]
                    if pd.notna(value):  # Skip NaN values
                        row_content.append(f"{col}: {value}")

                if row_content:
                    content = " | ".join(row_content)

                    # Base metadata
                    metadata = {
                        "filename": filename,
                        "file_type": "csv",
                        "row_number": row_idx + 2,  # +2 because pandas 0-indexed + header row
                        "column_headers": column_headers,
                        "upload_date": upload_date,
                    }

                    # Add temporal context if available
                    if row_idx in temporal_data_by_row and temporal_data_by_row[row_idx]:
                        metadata["temporal_context"] = temporal_data_by_row[row_idx]

                    chunk = DocumentChunk(
                        content=content,
                        metadata=metadata,
                        file_type=FileType.CSV,
                    )
                    chunks.append(chunk)

            logger.info(f"Parsed CSV file: {filename} - {len(chunks)} chunks")
            return chunks

        except Exception as e:
            logger.error(f"Error parsing CSV file {filename}: {e}")
            raise


class TextParser(DocumentParser):
    """Parser for plain text files with line-based chunking."""

    def parse(self, file_bytes: bytes, filename: str) -> List[DocumentChunk]:
        """
        Parse text file line-by-line or paragraph-by-paragraph.

        Returns chunks with metadata:
        - filename, file_type, line_number or paragraph_index
        """
        chunks = []
        upload_date = datetime.utcnow().isoformat()

        try:
            text = file_bytes.decode('utf-8', errors='ignore')

            # Split by double newlines (paragraphs)
            paragraphs = text.split('\n\n')

            for para_idx, para in enumerate(paragraphs):
                para = para.strip()
                if para:
                    # If paragraph too long, split further
                    if len(para) > 4000:
                        # Split by single newlines
                        lines = para.split('\n')
                        current_chunk = ""
                        line_start = 0

                        for line in lines:
                            if len(current_chunk) + len(line) < 4000:
                                current_chunk += line + "\n"
                            else:
                                if current_chunk.strip():
                                    chunk = DocumentChunk(
                                        content=current_chunk.strip(),
                                        metadata={
                                            "filename": filename,
                                            "file_type": "text",
                                            "paragraph_index": para_idx,
                                            "upload_date": upload_date,
                                        },
                                        file_type=FileType.TEXT,
                                    )
                                    chunks.append(chunk)
                                current_chunk = line + "\n"

                        if current_chunk.strip():
                            chunk = DocumentChunk(
                                content=current_chunk.strip(),
                                metadata={
                                    "filename": filename,
                                    "file_type": "text",
                                    "paragraph_index": para_idx,
                                },
                                file_type=FileType.TEXT,
                            )
                            chunks.append(chunk)
                    else:
                        chunk = DocumentChunk(
                            content=para,
                            metadata={
                                "filename": filename,
                                "file_type": "text",
                                "paragraph_index": para_idx,
                                "upload_date": upload_date,
                            },
                            file_type=FileType.TEXT,
                        )
                        chunks.append(chunk)

            logger.info(f"Parsed text file: {filename} - {len(chunks)} chunks")
            return chunks

        except Exception as e:
            logger.error(f"Error parsing text file {filename}: {e}")
            raise


class DocumentParserFactory:
    """Factory to get appropriate parser for file type."""

    _parsers = {
        FileType.EXCEL: ExcelParser(),
        FileType.CSV: CSVParser(),
        FileType.PDF: PDFParser(),
        FileType.WORD: WordParser(),
        FileType.POWERPOINT: PowerPointParser(),
        FileType.TEXT: TextParser(),
    }

    @classmethod
    def get_parser(cls, file_type: FileType) -> DocumentParser:
        """Get parser for given file type."""
        parser = cls._parsers.get(file_type)
        if not parser:
            raise ValueError(f"No parser available for file type: {file_type}")
        return parser
